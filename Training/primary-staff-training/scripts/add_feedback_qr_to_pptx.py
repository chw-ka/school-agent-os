#!/usr/bin/env python3
"""Add feedback form QR code slide to STEAM-Sharing-M5StickS3.pptx."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "docs" / "STEAM-Sharing-M5StickS3.pptx"
META = ROOT / "docs" / "feedback-form-meta.json"
QR_PNG = ROOT / "docs" / "assets" / "feedback-form-qr.png"


def load_form_url() -> str:
    if not META.exists():
        print(f"Run create_feedback_form.py first. Missing: {META}", file=sys.stderr)
        sys.exit(1)
    data = json.loads(META.read_text(encoding="utf-8"))
    url = data.get("responder_uri")
    if not url:
        print("No responder_uri in metadata.", file=sys.stderr)
        sys.exit(1)
    return url


def generate_qr(url: str) -> Path:
    QR_PNG.parent.mkdir(parents=True, exist_ok=True)
    img = qrcode.make(url, box_size=8, border=2)
    img.save(str(QR_PNG))
    print(f"QR saved: {QR_PNG}")
    return QR_PNG


def add_feedback_slide(prs: Presentation, url: str, qr_path: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    tf.text = "課後 Feedback — 請掃 QR Code"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x3A, 0x5C)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8))
    stf = sub_box.text_frame
    stf.text = "STEAM 分享：M5StickS3 · 2026 年 6 月 22 日 · 約 3–5 分鐘"
    sp = stf.paragraphs[0]
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    sp.alignment = PP_ALIGN.CENTER

    # QR code (center-left)
    slide.shapes.add_picture(str(qr_path), Inches(1.2), Inches(2.0), width=Inches(4.2))

    # Instructions (right)
    body_box = slide.shapes.add_textbox(Inches(5.8), Inches(2.2), Inches(6.8), Inches(4.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    lines = [
        "用手機掃描左邊 QR Code",
        "填寫工作坊 feedback form",
        "",
        "或喺 Google Classroom",
        "搵「課後 Feedback」連結",
        "",
        "多謝各位同事！",
        "有問題歡迎聯絡 Warren Chan",
    ]
    for i, line in enumerate(lines):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = line
        para.font.size = Pt(24 if i == 0 else 20)
        para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        if i == 0:
            para.font.bold = True

    # Short URL at bottom
    url_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.6))
    utf = url_box.text_frame
    utf.text = url
    up = utf.paragraphs[0]
    up.font.size = Pt(14)
    up.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    up.alignment = PP_ALIGN.CENTER


def main() -> None:
    if not PPTX.exists():
        print(f"Missing: {PPTX}", file=sys.stderr)
        sys.exit(1)

    url = load_form_url()
    qr_path = generate_qr(url)

    prs = Presentation(str(PPTX))
    last_title = ""
    if prs.slides:
        for shape in prs.slides[-1].shapes:
            if shape.has_text_frame and shape.text_frame.text:
                last_title = shape.text_frame.text.strip()
                break

    if "Feedback" in last_title and "QR" in last_title:
        print("Feedback slide already exists — replacing last slide.")
        r_id = prs.slides._sldIdLst[-1].rId  # noqa: SLF001
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[-1]

    add_feedback_slide(prs, url, qr_path)
    prs.save(str(PPTX))
    print(f"Updated: {PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
