#!/usr/bin/env python3
"""Build PowerPoint deck from docs/ppt-output PNGs + screenshot-embedded slides."""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parent.parent
OUTPUT_DIR = ROOT / "docs" / "ppt-output"
PROMPTS_MD = ROOT / "docs" / "ppt-slides-gemini-prompts.md"
OUT_PPTX = ROOT / "docs" / "STEAM-Sharing-M5StickS3.pptx"
CROPPED_DIR = ROOT / "docs" / "screen-captures" / "cropped"
EXTRACTED_DIR = ROOT / "docs" / "screen-captures" / "extracted"
ASSETS_DIR = ROOT / "docs" / "assets"

DARK_BLUE = RGBColor(0x1A, 0x3A, 0x5C)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY = RGBColor(0x66, 0x66, 0x66)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)


def slide_files() -> list[tuple[int, str]]:
    text = PROMPTS_MD.read_text(encoding="utf-8")
    rows: list[tuple[int, str]] = []
    for m in re.finditer(r"\|\s*(\d{2})\s*\|\s*`(slide-\d{2}-[^`]+\.png)`", text):
        rows.append((int(m.group(1)), m.group(2)))
    return sorted(rows, key=lambda x: x[0])


def add_title(slide, text: str, y=Inches(0.35), h=Inches(0.85),
              size=Pt(38), color=None, bold=True) -> None:
    color = color or DARK_BLUE
    tb = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.2), h)
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color


def add_text_box(slide, text: str, x, y, w, h,
                 size=Pt(18), color=None, bold=False, align=PP_ALIGN.LEFT,
                 wrap=True) -> None:
    color = color or DARK_GREY
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = align
        para.font.size = size
        para.font.bold = bold
        para.font.color.rgb = color


def place_image(slide, path: Path, x, y, w=None, h=None,
                max_w=None, max_h=None) -> None:
    """Add image with optional max-dimension scaling."""
    from PIL import Image as PILImage
    img = PILImage.open(path)
    iw, ih = img.size
    if w is None and h is None:
        if max_w and max_h:
            scale = min(max_w / iw, max_h / ih)
            w = int(iw * scale)
            h = int(ih * scale)
        elif max_w:
            scale = max_w / iw
            w = max_w
            h = int(ih * scale)
        elif max_h:
            scale = max_h / ih
            h = max_h
            w = int(iw * scale)
        else:
            w = iw
            h = ih
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


# ─── Slide 4: Know These Three Buttons ───────────────────────────────────────

def add_slide_04(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Light background
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    # Title bar
    title_bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.05))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()

    add_title(slide, "Know These Three Buttons", y=Inches(0.15),
              h=Inches(0.8), size=Pt(40), color=WHITE)

    # Device photo (reference image) centred on left half
    ref = ASSETS_DIR / "m5sticks3-product-reference.png"
    if ref.exists():
        place_image(slide, ref, Inches(0.4), Inches(1.2),
                    max_w=Inches(5.5), max_h=Inches(5.8))

    # Button labels on the right
    labels = [
        ("Power Button", "Hold 2 sec to power ON\nPress to wake / sleep", Inches(3.0)),
        ("Button A", "Left front button\n(BtnA in UIFlow code)", Inches(4.5)),
        ("Button B", "Right side button\n(BtnB in UIFlow code)", Inches(6.0)),
    ]
    for title_txt, body_txt, y in labels:
        # coloured pill
        pill = slide.shapes.add_shape(1, Inches(6.2), y, Inches(6.5), Inches(0.55))
        pill.fill.solid()
        pill.fill.fore_color.rgb = TEAL
        pill.line.fill.background()
        tf = pill.text_frame
        tf.text = title_txt
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        add_text_box(slide, body_txt, Inches(6.4), y + Inches(0.6),
                     Inches(6.2), Inches(0.8), size=Pt(17), color=DARK_GREY)


# ─── Slide 11: Live Demo xiaozhi.me ──────────────────────────────────────────

def add_slide_11(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    add_title(slide, "Live Demo: xiaozhi.me", y=Inches(0.5),
              h=Inches(0.9), size=Pt(44), color=WHITE)

    # URL badge
    badge = slide.shapes.add_shape(1, Inches(3.8), Inches(1.6), Inches(5.7), Inches(0.75))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = "xiaozhi.me"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE

    add_text_box(slide,
                 "Open browser → go to xiaozhi.me → log in\n\n"
                 "Show the Xiaozhi device management dashboard\n\n"
                 "Point out: Knowledge Base · MCP · Voice model settings",
                 Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.5),
                 size=Pt(22), color=RGBColor(0xCC, 0xE8, 0xFF),
                 align=PP_ALIGN.CENTER)

    # UIFlow IDE screenshot as footer reference
    sc = EXTRACTED_DIR / "screen-11.png"
    if sc.exists():
        place_image(slide, sc, Inches(0.5), Inches(5.5),
                    max_w=Inches(5.5), max_h=Inches(1.8))
        add_text_box(slide, "UIFlow 2 IDE — used in Task 2 & 3",
                     Inches(6.2), Inches(5.7), Inches(6.5), Inches(1.0),
                     size=Pt(16), color=RGBColor(0xAA, 0xCC, 0xFF))


# ─── Slide 12: Task 2 — OTA install UIFlow ────────────────────────────────────

def add_slide_12(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    title_bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.05))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()

    add_title(slide, "Task 2: Install UIFlow 2", y=Inches(0.15),
              h=Inches(0.8), size=Pt(38), color=WHITE)

    # Steps on left
    steps = [
        "① Power ON → M5Launcher appears",
        "② Press Button B immediately",
        "③ Select OTA in menu",
        "④ Choose UIFlow 2 from list",
        "⑤ Wait for install to complete",
    ]
    add_text_box(slide, "\n".join(steps),
                 Inches(0.45), Inches(1.25), Inches(5.4), Inches(5.0),
                 size=Pt(20), color=DARK_BLUE, bold=False)

    # Device photos on right (2×2 grid + one wide)
    photos = [
        ("d01-launcher-main-menu.jpg", Inches(6.0), Inches(1.2), Inches(3.4), Inches(1.9)),
        ("d02-ota-app-list.jpg",       Inches(9.6), Inches(1.2), Inches(3.4), Inches(1.9)),
        ("d03-uiflow-detail.jpg",      Inches(6.0), Inches(3.3), Inches(3.4), Inches(1.9)),
        ("d05-installing-fw.jpg",      Inches(9.6), Inches(3.3), Inches(3.4), Inches(1.9)),
    ]
    captions = ["① M5Launcher", "③ OTA list", "④ UIFlow 2", "⑤ Installing…"]
    for (fname, x, y, w, h), cap in zip(photos, captions):
        path = CROPPED_DIR / fname
        if path.exists():
            place_image(slide, path, x, y, max_w=w, max_h=h)
        add_text_box(slide, cap, x, y + h + Inches(0.05), w, Inches(0.35),
                     size=Pt(13), color=MID_GREY, align=PP_ALIGN.CENTER)

    # Badge
    badge = slide.shapes.add_shape(1, Inches(0.45), Inches(6.1), Inches(5.0), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = "⏱  Takes about 1–2 minutes"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE


# ─── Slide 14: Pair with UIFlow 2 ─────────────────────────────────────────────

def add_slide_14(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    title_bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.05))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = DARK_BLUE
    title_bg.line.fill.background()

    add_title(slide, "Pair with UIFlow 2", y=Inches(0.15),
              h=Inches(0.8), size=Pt(38), color=WHITE)

    # Steps on left
    steps = [
        "① On device: press Cloud icon",
        "② Note the 6-digit Access Code",
        "③ Open uiflow2.m5stack.com",
        "④ Click Connect Device → enter code",
        "⑤ Click Run Once to test",
    ]
    add_text_box(slide, "\n".join(steps),
                 Inches(0.45), Inches(1.25), Inches(5.5), Inches(4.5),
                 size=Pt(20), color=DARK_BLUE)

    # Top-right: device showing access code
    d07 = CROPPED_DIR / "d07-cloud-ok.jpg"
    if d07.exists():
        place_image(slide, d07, Inches(6.1), Inches(1.2),
                    max_w=Inches(3.0), max_h=Inches(4.5))
        add_text_box(slide, "② Device shows Access Code",
                     Inches(6.1), Inches(5.7), Inches(3.0), Inches(0.5),
                     size=Pt(14), color=MID_GREY, align=PP_ALIGN.CENTER)

    # Right: UIFlow website screenshot (Select Device dialog)
    sc6 = EXTRACTED_DIR / "screen-06.png"
    if sc6.exists():
        place_image(slide, sc6, Inches(9.3), Inches(1.2),
                    max_w=Inches(3.8), max_h=Inches(4.5))
        add_text_box(slide, "④ UIFlow 2: Connect Device",
                     Inches(9.3), Inches(5.7), Inches(3.8), Inches(0.5),
                     size=Pt(14), color=MID_GREY, align=PP_ALIGN.CENTER)

    # URL badge at bottom
    badge = slide.shapes.add_shape(1, Inches(0.45), Inches(6.1), Inches(5.2), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = "uiflow2.m5stack.com"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


# ─── Image slide ──────────────────────────────────────────────────────────────

def add_image_slide(prs: Presentation, image_path: Path) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(
        str(image_path), Inches(0), Inches(0),
        width=prs.slide_width, height=prs.slide_height,
    )


# ─── Blank fallback ───────────────────────────────────────────────────────────

def add_blank_slide(prs: Presentation, title: str, body: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
    btf = body_box.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body.split("\n")):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = line
        para.font.size = Pt(20)
        para.font.color.rgb = DARK_GREY


# ─── Dispatch map ─────────────────────────────────────────────────────────────

CUSTOM_SLIDES: dict[int, callable] = {
    4:  add_slide_04,
    11: add_slide_11,
    12: add_slide_12,
    14: add_slide_14,
}


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    missing: list[str] = []
    for num, filename in slide_files():
        if num in CUSTOM_SLIDES:
            CUSTOM_SLIDES[num](prs)
            print(f"[{num:02d}] custom slide built")
            continue

        path = OUTPUT_DIR / filename
        if not path.exists():
            missing.append(filename)
            add_blank_slide(
                prs,
                f"Slide {num} — missing image",
                f"Expected: {filename}\nRegenerate with generate_ppt_slides.py",
            )
            continue

        add_image_slide(prs, path)
        print(f"[{num:02d}] {filename}")

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"\nSaved: {OUT_PPTX} ({len(prs.slides)} slides)")
    if missing:
        print("Warning: missing PNGs:", ", ".join(missing))


if __name__ == "__main__":
    main()
